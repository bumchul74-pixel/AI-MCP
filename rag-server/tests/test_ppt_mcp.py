from __future__ import annotations

import os
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation

from app.ppt.ppt_server import (
    _analyze_template,
    _generation_error,
    _resolve_template_path,
    _safe_output_path,
    _validate_generation_request,
)


class PptMcpTest(unittest.TestCase):
    def _create_template(self, directory: Path) -> Path:
        template = directory / "template.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "Template title"
        presentation.save(template)
        return template

    def test_analyzes_template_layouts_and_source_slides(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            template = self._create_template(Path(directory))
            result = _analyze_template(template)

        self.assertGreater(result["layout_count"], 0)
        self.assertEqual(1, result["source_slide_count"])
        self.assertEqual("Template title", result["source_slides"][0]["title"])
        self.assertIn("title", result["layouts"][0]["capabilities"])

    def test_template_access_is_restricted_to_allowed_directories(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            allowed = Path(directory).resolve()
            template = self._create_template(allowed)
            with patch.dict(os.environ, {"PPT_MCP_ALLOWED_DIRS": str(allowed)}):
                self.assertEqual(template, _resolve_template_path(str(template)))

            outside = allowed.parent / "outside.pptx"
            with patch.dict(os.environ, {"PPT_MCP_ALLOWED_DIRS": str(allowed)}):
                with self.assertRaisesRegex(ValueError, "outside"):
                    _resolve_template_path(str(outside))

    def test_output_name_is_sanitized_and_confined(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            output_dir = Path(directory).resolve()
            with patch.dict(os.environ, {"PPT_MCP_OUTPUT_DIR": str(output_dir)}):
                output = _safe_output_path("../분기 보고서?.pptx", "ignored")

        self.assertEqual(output_dir, output.parent)
        self.assertEqual("분기_보고서.pptx", output.name)

    def test_generation_request_validation(self) -> None:
        _validate_generation_request("content", 10, "Korean", "mock", "corporate")
        with self.assertRaisesRegex(ValueError, "between 3 and 50"):
            _validate_generation_request("content", 2, "Korean", "mock", "corporate")
        with self.assertRaisesRegex(ValueError, "provider"):
            _validate_generation_request("content", 10, "Korean", "unknown", "corporate")

    def test_generation_error_classifies_provider_quota_variants(self) -> None:
        expected = (
            "AutoPPT provider rate limit was exceeded. "
            "Retry later or choose another provider."
        )
        self.assertEqual(expected, _generation_error("HTTP 429 Too Many Requests", ""))
        self.assertEqual(expected, _generation_error("insufficient_quota", ""))


if __name__ == "__main__":
    unittest.main()
