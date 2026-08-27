from __future__ import annotations

import hashlib
import json
import logging
import os
import re
import subprocess
import sys
import tempfile
import zipfile
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path
from threading import BoundedSemaphore
from typing import Any

from mcp.server.fastmcp import FastMCP
from pptx import Presentation


logger = logging.getLogger(__name__)
PPT_SERVER_DIR = Path(__file__).resolve().parents[2]
PROJECT_DIR = PPT_SERVER_DIR.parent
DEFAULT_ALLOWED_DIRS = (PPT_SERVER_DIR / "inbox", PROJECT_DIR / "uploads")
DEFAULT_OUTPUT_DIR = PPT_SERVER_DIR / "output" / "ppt"
JOB_DIR = PPT_SERVER_DIR / "data" / "ppt-jobs"
SUPPORTED_PROVIDERS = {"openai", "google", "anthropic", "mock"}
SUPPORTED_STYLES = {
    "minimalist", "technology", "nature", "creative", "corporate", "academic",
    "startup", "dark", "luxury", "magazine", "tech_gradient", "ocean", "sunset",
    "chalkboard", "blueprint", "sketch", "retro", "neon",
}
MAX_TEMPLATE_BYTES = 100 * 1024 * 1024
MAX_DECOMPRESSED_BYTES = 500 * 1024 * 1024
MAX_CONTENT_CHARS = 50_000
_generation_slots = BoundedSemaphore(max(1, int(os.getenv("PPT_MCP_MAX_CONCURRENT", "1"))))


def _configured_directories(name: str, defaults: tuple[Path, ...]) -> tuple[Path, ...]:
    configured = os.getenv(name)
    values = configured.split(os.pathsep) if configured else defaults
    return tuple(Path(value).expanduser().resolve() for value in values if str(value).strip())


def _allowed_directories() -> tuple[Path, ...]:
    return _configured_directories("PPT_MCP_ALLOWED_DIRS", DEFAULT_ALLOWED_DIRS)


def _output_directory() -> Path:
    configured = os.getenv("PPT_MCP_OUTPUT_DIR")
    output_dir = Path(configured).expanduser().resolve() if configured else DEFAULT_OUTPUT_DIR
    output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir


def _is_within(path: Path, directories: tuple[Path, ...]) -> bool:
    return any(path == directory or path.is_relative_to(directory) for directory in directories)


def _validate_pptx_archive(path: Path) -> None:
    if path.stat().st_size > MAX_TEMPLATE_BYTES:
        raise ValueError("PPTX file exceeds the 100 MB size limit.")
    try:
        with zipfile.ZipFile(path) as archive:
            total_size = 0
            for entry in archive.infolist():
                parts = entry.filename.replace("\\", "/").split("/")
                if entry.filename.startswith("/") or ".." in parts:
                    raise ValueError("PPTX contains an unsafe archive entry.")
                total_size += entry.file_size
                if total_size > MAX_DECOMPRESSED_BYTES:
                    raise ValueError("PPTX expanded content exceeds the 500 MB size limit.")
            if "[Content_Types].xml" not in archive.namelist():
                raise ValueError("File is not a valid PPTX package.")
    except zipfile.BadZipFile as exception:
        raise ValueError("File is not a valid PPTX package.") from exception


def _resolve_template_path(file_path: str) -> Path:
    if not file_path or not file_path.strip():
        raise ValueError("file_path is required.")
    candidate = Path(file_path.strip()).expanduser()
    path = candidate.resolve() if candidate.is_absolute() else (PPT_SERVER_DIR / candidate).resolve()
    if not _is_within(path, _allowed_directories()):
        raise ValueError("Template path is outside PPT_MCP_ALLOWED_DIRS.")
    if path.suffix.lower() != ".pptx":
        raise ValueError("Only .pptx template files are supported.")
    if not path.is_file():
        raise ValueError(f"Template file was not found: {path}")
    _validate_pptx_archive(path)
    return path


def _safe_output_path(output_name: str, content: str) -> Path:
    requested = Path(output_name).name if output_name and output_name.strip() else ""
    stem_source = Path(requested).stem if requested else content.splitlines()[0][:80]
    safe_stem = re.sub(r"[^0-9A-Za-z가-힣_-]+", "_", stem_source).strip("_.-")
    safe_stem = safe_stem[:100] or "generated_presentation"
    output_dir = _output_directory()
    target = (output_dir / f"{safe_stem}.pptx").resolve()
    sequence = 2
    while target.exists():
        target = (output_dir / f"{safe_stem}-{sequence}.pptx").resolve()
        sequence += 1
    return target


def _placeholder_role(type_name: str) -> str:
    normalized = type_name.upper()
    if "TITLE" in normalized:
        return "title"
    if "PICTURE" in normalized or "BITMAP" in normalized:
        return "image"
    if any(value in normalized for value in ("BODY", "OBJECT", "TEXT")):
        return "content"
    return "other"


def _layout_capabilities(layout_name: str, roles: list[str]) -> list[str]:
    name = layout_name.lower()
    capabilities: set[str] = set()
    if "title" in roles:
        capabilities.add("title")
    content_count = roles.count("content")
    if content_count:
        capabilities.add("content")
    if content_count >= 2 or any(token in name for token in ("two", "comparison", "2개", "비교")):
        capabilities.update(("two_column", "comparison"))
    if "image" in roles or any(token in name for token in ("picture", "image", "사진", "그림")):
        capabilities.add("image")
    if any(token in name for token in ("section", "header", "구역", "섹션")):
        capabilities.add("section")
    if not roles or "blank" in name or "빈 화면" in name:
        capabilities.add("blank")
    return sorted(capabilities)


def _analyze_template(path: Path) -> dict[str, Any]:
    presentation = Presentation(str(path))
    layouts: list[dict[str, Any]] = []
    flat_index = 0
    for master_index, master in enumerate(presentation.slide_masters):
        for layout_index, layout in enumerate(master.slide_layouts):
            placeholders = []
            roles = []
            for placeholder in layout.placeholders:
                type_name = str(placeholder.placeholder_format.type)
                role = _placeholder_role(type_name)
                roles.append(role)
                placeholders.append({
                    "index": placeholder.placeholder_format.idx,
                    "name": placeholder.name,
                    "type": type_name,
                    "role": role,
                })
            layouts.append({
                "index": flat_index,
                "master_index": master_index,
                "layout_index": layout_index,
                "name": layout.name,
                "capabilities": _layout_capabilities(layout.name, roles),
                "placeholders": placeholders,
            })
            flat_index += 1

    source_slides = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_values = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        source_slides.append({
            "slide_number": slide_number,
            "layout_name": slide.slide_layout.name,
            "title": slide.shapes.title.text.strip() if slide.shapes.title else "",
            "text_preview": " | ".join(text_values)[:500],
            "shape_count": len(slide.shapes),
        })

    return {
        "template_path": str(path),
        "file_size": path.stat().st_size,
        "slide_width_inches": round(presentation.slide_width / 914400, 3),
        "slide_height_inches": round(presentation.slide_height / 914400, 3),
        "master_count": len(presentation.slide_masters),
        "layout_count": len(layouts),
        "source_slide_count": len(source_slides),
        "layouts": layouts,
        "source_slides": source_slides,
        "engine": "AutoPPT",
    }


def _package_version(package: str) -> str:
    try:
        return version(package)
    except PackageNotFoundError:
        return "unknown"


def _validate_generation_request(
        content: str, slides: int, language: str, provider: str, style: str) -> None:
    if not content or not content.strip():
        raise ValueError("content is required.")
    if len(content) > MAX_CONTENT_CHARS:
        raise ValueError(f"content must not exceed {MAX_CONTENT_CHARS} characters.")
    if not 3 <= slides <= 50:
        raise ValueError("slides must be between 3 and 50.")
    if not language or len(language) > 50:
        raise ValueError("language is required and must not exceed 50 characters.")
    if provider not in SUPPORTED_PROVIDERS:
        raise ValueError(f"provider must be one of: {', '.join(sorted(SUPPORTED_PROVIDERS))}.")
    if style not in SUPPORTED_STYLES:
        raise ValueError(f"style must be one of: {', '.join(sorted(SUPPORTED_STYLES))}.")


def _generation_error(stderr: str, stdout: str) -> str:
    combined = f"{stderr}\n{stdout}".lower()
    if "api key" in combined or "authentication" in combined:
        return "AutoPPT authentication failed. Check the selected provider API key."
    if any(signal in combined for signal in (
            "rate limit", "too many requests", "insufficient_quota", "quota", "429")):
        return "AutoPPT provider rate limit was exceeded. Retry later or choose another provider."
    if "template" in combined and "not found" in combined:
        return "AutoPPT could not read the selected template."
    return "AutoPPT generation failed. Check provider configuration and the PPT MCP server log."


def _run_generation(payload: dict[str, Any], timeout_seconds: int) -> None:
    JOB_DIR.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="ppt-", dir=JOB_DIR) as job_directory:
        request_path = Path(job_directory) / "request.json"
        request_path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
        result = subprocess.run(
            [sys.executable, "-m", "app.ppt.autoppt_runner", str(request_path)],
            cwd=PPT_SERVER_DIR,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout_seconds,
            check=False,
        )
    if result.returncode != 0:
        safe_error = _generation_error(result.stderr, result.stdout)
        logger.warning("AutoPPT child process failed. category=%s", safe_error)
        raise RuntimeError(safe_error)


ppt_server = FastMCP(
    "ppt-generator",
    instructions=(
        "Analyze an uploaded PPTX template and generate an editable PowerPoint deck with AutoPPT. "
        "Template and output access is restricted to configured directories."
    ),
    host=os.getenv("PPT_MCP_HOST", "127.0.0.1"),
    port=int(os.getenv("PPT_MCP_PORT", "8002")),
    streamable_http_path="/ppt",
)


@ppt_server.tool(description=(
    "Analyze a PPTX template's masters, layouts, placeholders, dimensions, and source slides "
    "before generation. The file must be inside PPT_MCP_ALLOWED_DIRS."
))
def analyze_ppt_template(file_path: str) -> dict[str, Any]:
    return _analyze_template(_resolve_template_path(file_path))


@ppt_server.tool(description=(
    "Generate an editable PPTX from content using the selected template. AutoPPT plans the deck, "
    "selects content-appropriate slide types, preserves the template package, and performs deck QA."
))
def generate_presentation(
        content: str,
        template_path: str,
        output_name: str = "",
        slides: int = 10,
        language: str = "Korean",
        provider: str = "",
        model: str = "",
        style: str = "corporate",
        offline: bool = False,
        thumbnails: bool = False) -> dict[str, Any]:
    selected_provider = ("mock" if offline else provider.strip().lower()) or os.getenv(
        "PPT_MCP_PROVIDER", "openai").strip().lower()
    selected_style = style.strip().lower() or "corporate"
    _validate_generation_request(content, slides, language, selected_provider, selected_style)
    resolved_template = _resolve_template_path(template_path)
    output_path = _safe_output_path(output_name, content)
    timeout_seconds = max(30, int(os.getenv("PPT_MCP_GENERATION_TIMEOUT_SECONDS", "900")))
    payload = {
        "content": content.strip(),
        "template_path": str(resolved_template),
        "output_path": str(output_path),
        "slides": slides,
        "language": language.strip(),
        "provider": selected_provider,
        "model": model.strip() or os.getenv("PPT_MCP_MODEL") or None,
        "style": selected_style,
        "offline": offline,
        "thumbnails": thumbnails,
    }
    try:
        with _generation_slots:
            _run_generation(payload, timeout_seconds)
    except subprocess.TimeoutExpired as exception:
        raise RuntimeError(f"PPT generation exceeded the {timeout_seconds}-second limit.") from exception

    if not output_path.is_file():
        raise RuntimeError("AutoPPT completed without producing the expected PPTX file.")
    _validate_pptx_archive(output_path)
    digest = hashlib.sha256(output_path.read_bytes()).hexdigest()
    template_analysis = _analyze_template(resolved_template)
    generated_analysis = _analyze_template(output_path)
    return {
        "output_path": str(output_path),
        "file_name": output_path.name,
        "file_size": output_path.stat().st_size,
        "sha256": digest,
        "slides_requested": slides,
        "language": language.strip(),
        "provider": selected_provider,
        "template_path": str(resolved_template),
        "template_layout_count": template_analysis["layout_count"],
        "slides_generated": generated_analysis["source_slide_count"],
        "engine": "AutoPPT",
        "engine_version": _package_version("autoppt"),
    }


def main() -> None:
    transport = os.getenv("PPT_MCP_TRANSPORT", "stdio").strip().lower()
    if transport not in {"stdio", "sse", "streamable-http"}:
        raise ValueError("PPT_MCP_TRANSPORT must be stdio, sse, or streamable-http.")
    ppt_server.run(transport=transport)


if __name__ == "__main__":
    main()
